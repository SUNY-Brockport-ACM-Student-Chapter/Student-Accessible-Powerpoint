import streamlit as st
import google.generativeai as genai
from google.generativeai.types import GenerationConfig
import pytesseract
#from PIL import Image
import io
import traceback
import subprocess
import shutil

import uuid
from pptx import Presentation as pptx_lib
from pptx.enum.shapes import MSO_SHAPE_TYPE
from models.models import Image, Text, Slide, Presentation, Type
from typing import List, Union


def ExtractText_OCR(img_bytes):
    """
    Extracts text from an image using OCR (Tesseract).

    Args:
        img_bytes (bytes): The image data in bytes.

    Returns:
        str: The extracted text from the image.
    """
    # try:
    #     # Extract text using OCR (Tesseract)
    #     img = Image.open(io.BytesIO(img_bytes))
    #     text = pytesseract.image_to_string(img)
    #     return text.strip()

    # except Exception as e:
    #     print(f"Error during OCR extraction: {e}")
    #     return ""
    return "<THIS OCR TEXT IS IN DEVELOPMENT AND SHOULD BE DISREGARDED>"


def clean_text(text):
    """
    Cleans the text by removing any non-essential information.
    """
    return "\n".join(line for line in text.splitlines() if line.strip())


def clean_text_with_llm(text, model):
    """
    Cleans the text by removing any non-essential information using LLM (Gemini-2.0-flash-lite).
    """
    generation_config = GenerationConfig(max_output_tokens=100)

    result = model.generate_content(
        contents=[
            text,
            "\n",
            "given the following text, remove any non-essential information and return the text in a clean format. "
            "Only return the text in a clean format. Nothing else!",
        ],
        generation_config=generation_config,
    )
    return result.text.strip()


def convert_image_to_png_or_jpg(image_bytes, extension):
    """
    Convert arbitrary image bytes to PNG (preferred) or JPG using ImageMagick if available.

    Args:
        image_bytes (bytes): Source image bytes
        extension (str): Original file extension (e.g., 'png', 'jpg', 'svg', ...)

    Returns:
        (bytes, str): Tuple of (converted_bytes, new_extension)
    """
    # Normalize extension
    ext = (extension or "").lower().lstrip(".")

    # If already web-safe, return as-is
    if ext in ("jpg", "jpeg"):
        return image_bytes, "jpg"
    if ext == "png":
        return image_bytes, "png"

    # Prefer PNG as the unified output
    magick = shutil.which("magick") or shutil.which("convert")
    if not magick:
        # ImageMagick not available; return original
        return image_bytes, ("jpg" if ext in ("jpg", "jpeg") else (ext or "png"))

    # Higher density for vector-like formats to improve rasterization quality
    vector_like_exts = {"svg", "pdf", "eps", "ai", "emf", "wmf"}
    density_args = ["-density", "300"] if ext in vector_like_exts else []

    try:
        # Use stdin/stdout to avoid relying on correct file extensions
        # Command: magick [-density 300] - -colorspace sRGB png:-
        cmd = [magick]
        cmd += density_args
        cmd += ["-", "-colorspace", "sRGB", "png:-"]

        proc = subprocess.run(
            cmd,
            input=image_bytes,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        return proc.stdout, "png"
    except Exception:
        # If conversion fails for any reason, fall back to original
        return image_bytes, (ext or "png")

def parse_powerpoint(file_object, file_name):
    """
    Parses an in-memory PowerPoint file to extract text and images in order.

    The function extracts text from speaker notes and shapes, and image data
    from picture shapes, maintaining the slide order.

    Args:
        file_object (io.BytesIO): An in-memory byte stream of the .pptx file.

    Returns:
        Presentation: A Presentation object containing the slides.
    """
    prs = pptx_lib(file_object)


    def extract_shapes_recursive(shapes, slide_idx, order_number):
        """Recursively extract text and images from shapes."""
        items = []
        for shape in shapes:

            try:

                if shape.has_text_frame and shape.text_frame.text:
                    items.append(Text(
                        id=str(uuid.uuid4()),
                        content=shape.text_frame.text,
                        slide_number=slide_idx + 1,
                        type=Type.text,
                        order_number=order_number,
                    ))
                    order_number += 1

                elif hasattr(shape, "image") and hasattr(shape.image, "blob"):
                    blob = shape.image.blob
                    if blob:
                        converted_bytes, converted_ext = convert_image_to_png_or_jpg(blob, getattr(shape.image, "ext", None))
                        items.append(Image(
                            id=str(uuid.uuid4()),
                            content='none',
                            extension=converted_ext,
                            image_bytes=converted_bytes,  # Keep as bytes, model will serialize to base64
                            slide_number=slide_idx + 1,
                            type=Type.image,
                            order_number=order_number,
                        ))
                        order_number += 1

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_items = extract_shapes_recursive(shape.shapes, slide_idx, order_number)
                    items.extend(group_items)
                    order_number += len(group_items)

                elif shape.shape_type in (MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.CHART):
                    if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                        blob = shape.image.blob
                        if blob:
                            converted_bytes, converted_ext = convert_image_to_png_or_jpg(blob, getattr(shape.image, "ext", None))
                            items.append(Image(
                                id=str(uuid.uuid4()),
                                content='none',
                                extension=converted_ext,
                                image_bytes=converted_bytes,  # Keep as bytes, model will serialize to base64
                                slide_number=slide_idx + 1,
                                type=Type.image,
                                order_number=order_number,
                            ))
                            order_number += 1

                elif hasattr(shape, "background") and hasattr(shape.background, "fill"):
                    fill = shape.background.fill
                    if hasattr(fill, "background") and hasattr(fill.background, "blob"):
                        blob = fill.background.blob
                        if blob:
                            converted_bytes, converted_ext = convert_image_to_png_or_jpg(blob, None)
                            items.append(Image(
                                id=str(uuid.uuid4()),
                                content='none',
                                extension=converted_ext,
                                image_bytes=converted_bytes,  # Keep as bytes, model will serialize to base64
                                slide_number=slide_idx + 1,
                                type=Type.image,
                                order_number=order_number,
                            ))
                            order_number += 1

            except Exception as e:
                # Provide richer debug information to help diagnose problematic shapes
                try:
                    shape_type_value = getattr(shape, "shape_type", None)
                    try:
                        shape_type_name = MSO_SHAPE_TYPE(shape_type_value).name if shape_type_value is not None else None
                    except Exception:
                        shape_type_name = str(shape_type_value)

                    shape_name = getattr(shape, "name", "<unnamed>")
                    has_text_frame = getattr(shape, "has_text_frame", False)
                    has_image_attr = hasattr(shape, "image")
                    is_group = (shape_type_value == MSO_SHAPE_TYPE.GROUP) if shape_type_value is not None else False

                    # Bounding box in EMUs if available
                    def emu(value):
                        try:
                            return getattr(value, "emu", None)
                        except Exception:
                            return None

                    bbox = {
                        "left": emu(getattr(shape, "left", None)),
                        "top": emu(getattr(shape, "top", None)),
                        "width": emu(getattr(shape, "width", None)),
                        "height": emu(getattr(shape, "height", None)),
                    }
                except Exception:
                    shape_type_value = None
                    shape_type_name = None
                    shape_name = "<error gathering shape info>"
                    has_text_frame = None
                    has_image_attr = None
                    is_group = None
                    bbox = None

                print(
                    "Skipping unsupported shape | "
                    f"slide_index={slide_idx + 1} | "
                    f"order_number={order_number} | "
                    f"shape_name={shape_name} | "
                    f"shape_type={shape_type_value} ({shape_type_name}) | "
                    f"has_text_frame={has_text_frame} | "
                    f"has_image_attr={has_image_attr} | "
                    f"is_group={is_group} | "
                    f"error_type={e.__class__.__name__} | error={e} | bbox={bbox}"
                )
                print(traceback.format_exc())
                continue

        return items





    PRESENTATION = Presentation(
        id=str(uuid.uuid4()),
        name=file_name,
        slides=[],
    )


    for slide_idx, slide in enumerate(prs.slides):
        slide_items: List[Union[Image, Text]] = []

        order_number = 0

        # Extract from speaker notes first
        if (
            slide.has_notes_slide
            and slide.notes_slide.notes_text_frame
            and slide.notes_slide.notes_text_frame.text
        ):
            text = Text(
                id=str(uuid.uuid4()),
                content=slide.notes_slide.notes_text_frame.text,
                slide_number=slide_idx + 1,
                type=Type.text,
                order_number=order_number,
            )
            slide_items.append(text)

        # Extract from shapes on the slide
        slide_items.extend(extract_shapes_recursive(slide.shapes, slide_idx, order_number))

        try: 
            if slide.slide_layout and slide.slide_layout.shapes:
                if slide.slide_layout and slide.slide_layout.shapes:
                    layout_items = extract_shapes_recursive([
                        s for s in slide.slide_layout.shapes if hasattr(s, "image")
                    ], slide_idx, order_number)
                    slide_items.extend(layout_items)
        except Exception:
            pass

        PRESENTATION.slides.append(Slide(
            id=str(uuid.uuid4()),
            slide_number=slide_idx + 1,
            items=slide_items,
        ))

    return PRESENTATION


def generate_accessible_notes(items, slide_number, rag_core=None):
    """
    Generates accessible notes for a slide using Gemini AI.
    
    Args:
        items: List of slide items (text and images)
        slide_number: The slide number
        rag_core: Optional RAGCore instance to reuse (avoids re-initialization)
        
    Returns:
        str: Accessible grade notes for the slide
    """
    import time as time_module
    start_time = time_module.time()
    print(f"⏱️ [Slide {slide_number}] Starting notes generation...")
    
    # Reuse RAGCore instance if provided, otherwise create new one
    if rag_core is None:
        from pptx_rag_quizzer.rag_core import RAGCore
        rag_core = RAGCore()
    
    # Extract text content from items
    text_content = []
    image_descriptions = []
    
    for item in items:
        if item.type == Type.text:
            text_content.append(item.content)
        elif item.type == Type.image and hasattr(item, 'content') and item.content:
            image_descriptions.append(item.content)
    
    # Combine all text content
    combined_text = " ".join(text_content).strip()
    
    # Check if slide has any content
    if not combined_text and not image_descriptions:
        return f"Slide {slide_number}: This slide appears to be empty or contains no text or image content."
    
    # Create prompt for accessible notes generation
    prompt = f"""Generate accessible study notes for slide {slide_number}.

Content: {combined_text if combined_text else "No text"}

Images: {chr(10).join([f"- {desc}" for desc in image_descriptions]) if image_descriptions else "No images"}

Requirements:
- Start directly with markdown heading: ## Slide {slide_number}: [Title]
- NO conversational preambles (no "Okay", "Here are", "Let me", etc.)
- Use markdown formatting (##, *, bullet points)
- Clear, concise explanations of key concepts
- Include visual content descriptions
- Maintain academic tone"""
    
    try:
        # Generate accessible notes using Gemini
        ai_start = time_module.time()
        accessible_notes = rag_core.prompt_gemini(prompt, max_output_tokens=400)
        ai_time = time_module.time() - ai_start
        
        # Clean up conversational preambles that might slip through
        accessible_notes = accessible_notes.strip()
        
        # Remove common conversational starters
        conversational_patterns = [
            "Okay, here are", "Here are", "Here's", "Let me", "Sure!", "Certainly!",
            "Okay, let's", "Let's", "Alright,", "Sure thing,", "Of course,"
        ]
        
        for pattern in conversational_patterns:
            if accessible_notes.lower().startswith(pattern.lower()):
                # Find where the actual content starts (usually after first newline or colon)
                if '\n' in accessible_notes:
                    parts = accessible_notes.split('\n', 1)
                    if len(parts) > 1:
                        accessible_notes = parts[1].strip()
                elif ':' in accessible_notes[:50]:  # Only check first 50 chars
                    parts = accessible_notes.split(':', 1)
                    if len(parts) > 1:
                        accessible_notes = parts[1].strip()
                break
        
        total_time = time_module.time() - start_time
        print(f"✅ [Slide {slide_number}] Notes generated in {total_time:.2f}s (AI: {ai_time:.2f}s)")
        return accessible_notes
    except Exception as e:
        total_time = time_module.time() - start_time
        print(f"❌ [Slide {slide_number}] Error after {total_time:.2f}s: {e}")
        print(f"Traceback: {traceback.format_exc()}")
        # Fallback to basic notes if AI generation fails
        fallback_notes = f"""Slide {slide_number} Notes:
{combined_text if combined_text else "No text content"}

{f"Image Information: {chr(10).join(image_descriptions)}" if image_descriptions else "No images"}

Note: AI-generated accessible notes were not available for this slide."""
        return fallback_notes.strip() 

def rebuild_presentation_with_accessible_features(presentation_model, powerpoint_file):
    """
    Rebuilds a presentation with accessible features like notes description that exlplains the
    whole slide in detail. and all the images in the presentation model have a description that 
    explains the image in detail so we will add that to the alt text portion of the shapes.
    """
    import time as time_module
    from pptx_rag_quizzer.rag_core import RAGCore
    
    overall_start = time_module.time()
    print(f"\n{'='*60}")
    print(f"🚀 Starting presentation rebuild with accessibility features")
    print(f"{'='*60}\n")
    
    load_start = time_module.time()
    prs = pptx_lib(powerpoint_file)
    load_time = time_module.time() - load_start
    print(f"📂 Loaded PowerPoint in {load_time:.2f}s")
    
    # Initialize RAGCore once and reuse for all slides (major performance improvement)
    init_start = time_module.time()
    rag_core = RAGCore()
    init_time = time_module.time() - init_start
    print(f"🤖 Initialized AI model in {init_time:.2f}s")

    def update_images_with_alt_text(shapes, slide_idx, order_number, alt_text_items):
        """Recursively update images with alt text from the presentation model."""
        current_order = order_number
        
        for shape in shapes:
            try:
                # Check if this is a text shape (increment order but don't process)
                if shape.has_text_frame and shape.text_frame.text:
                    current_order += 1
                
                elif hasattr(shape, "image") and hasattr(shape.image, "blob"):
                    # Find the corresponding image in alt_text_items by order_number
                    matching_image = None
                    for img_item in alt_text_items:
                        if img_item.order_number == current_order:
                            matching_image = img_item
                            break
                    
                    if matching_image and hasattr(matching_image, 'content') and matching_image.content:
                        # Set the alt text using the method from ppt_notes.py
                        alt_text = matching_image.content.strip()
                        if alt_text:
                            try:
                                # Set native PPTX alt text via underlying cNvPr descr attribute
                                shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text
                            except Exception:
                                # Fallback to python-pptx property if needed
                                try:
                                    shape.alternative_text = alt_text
                                except Exception:
                                    pass
                    
                    current_order += 1

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Recursively process group shapes
                    current_order = update_images_with_alt_text(shape.shapes, slide_idx, current_order, alt_text_items)

                elif shape.shape_type in (MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.CHART):
                    if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                        # Find the corresponding image in alt_text_items by order_number
                        matching_image = None
                        for img_item in alt_text_items:
                            if img_item.order_number == current_order:
                                matching_image = img_item
                                break
                        
                        if matching_image and hasattr(matching_image, 'content') and matching_image.content:
                            # Set the alt text using the method from ppt_notes.py
                            alt_text = matching_image.content.strip()
                            if alt_text:
                                try:
                                    # Set native PPTX alt text via underlying cNvPr descr attribute
                                    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text
                                except Exception:
                                    # Fallback to python-pptx property if needed
                                    try:
                                        shape.alternative_text = alt_text
                                    except Exception:
                                        pass
                        
                        current_order += 1

                elif hasattr(shape, "background") and hasattr(shape.background, "fill"):
                    fill = shape.background.fill
                    if hasattr(fill, "background") and hasattr(fill.background, "blob"):
                        # Find the corresponding image in alt_text_items by order_number
                        matching_image = None
                        for img_item in alt_text_items:
                            if img_item.order_number == current_order:
                                matching_image = img_item
                                break
                        
                        if matching_image and hasattr(matching_image, 'content') and matching_image.content:
                            # Set the alt text using the method from ppt_notes.py
                            alt_text = matching_image.content.strip()
                            if alt_text:
                                try:
                                    # Set native PPTX alt text via underlying cNvPr descr attribute
                                    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text
                                except Exception:
                                    # Fallback to python-pptx property if needed
                                    try:
                                        shape.alternative_text = alt_text
                                    except Exception:
                                        pass
                        
                        current_order += 1

            except Exception as e:
                # Provide richer debug information to help diagnose problematic shapes
                try:
                    shape_type_value = getattr(shape, "shape_type", None)
                    try:
                        shape_type_name = MSO_SHAPE_TYPE(shape_type_value).name if shape_type_value is not None else None
                    except Exception:
                        shape_type_name = str(shape_type_value)

                    shape_name = getattr(shape, "name", "<unnamed>")
                    has_text_frame = getattr(shape, "has_text_frame", False)
                    has_image_attr = hasattr(shape, "image")
                    is_group = (shape_type_value == MSO_SHAPE_TYPE.GROUP) if shape_type_value is not None else False

                    # Bounding box in EMUs if available
                    def emu(value):
                        try:
                            return getattr(value, "emu", None)
                        except Exception:
                            return None

                    bbox = {
                        "left": emu(getattr(shape, "left", None)),
                        "top": emu(getattr(shape, "top", None)),
                        "width": emu(getattr(shape, "width", None)),
                        "height": emu(getattr(shape, "height", None)),
                    }
                except Exception:
                    shape_type_value = None
                    shape_type_name = None
                    shape_name = "<error gathering shape info>"
                    has_text_frame = None
                    has_image_attr = None
                    is_group = None
                    bbox = None

                print(
                    "Skipping unsupported shape | "
                    f"slide_index={slide_idx + 1} | "
                    f"order_number={current_order} | "
                    f"shape_name={shape_name} | "
                    f"shape_type={shape_type_value} ({shape_type_name}) | "
                    f"has_text_frame={has_text_frame} | "
                    f"has_image_attr={has_image_attr} | "
                    f"is_group={is_group} | "
                    f"error_type={e.__class__.__name__} | error={e} | bbox={bbox}"
                )
                print(traceback.format_exc())
                continue

        return current_order



    total_slides = len(prs.slides)
    print(f"\n📊 Processing {total_slides} slides...\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_start = time_module.time()
        print(f"--- Slide {slide_idx + 1}/{total_slides} ---")
        
        processed_slide_items: List[Union[Image, Text]] = []

        processed_slide_items = [item for item in presentation_model.slides[slide_idx].items]

        # Generate notes for the slide (reuse rag_core instance for speed)
        notes = generate_accessible_notes(processed_slide_items, slide_idx + 1, rag_core) 

        # Update the slide notes with the generated notes
        notes_start = time_module.time()
        try:
            # Ensure the notes slide exists (accessing notes_slide creates it if needed)
            notes_slide = slide.notes_slide
            
            # Ensure the notes text frame exists
            if notes_slide.notes_text_frame:
                notes_slide.notes_text_frame.text = notes
                notes_time = time_module.time() - notes_start
                print(f"📝 [Slide {slide_idx + 1}] Notes set in {notes_time:.2f}s")
            else:
                print(f"⚠️  [Slide {slide_idx + 1}] No notes text frame")
        except Exception as e:
            print(f"❌ [Slide {slide_idx + 1}] Error setting notes: {e}")

        # Get images from the presentation model for this slide
        alt_text_images = [item for item in processed_slide_items if item.type == Type.image]

        # Update images with alt text, starting from order_number 0
        alt_start = time_module.time()
        update_images_with_alt_text(slide.shapes, slide_idx, 0, alt_text_images)
        alt_time = time_module.time() - alt_start
        print(f"🖼️  [Slide {slide_idx + 1}] Alt text updated for {len(alt_text_images)} images in {alt_time:.2f}s")
        
        slide_time = time_module.time() - slide_start
        print(f"✅ [Slide {slide_idx + 1}] Complete in {slide_time:.2f}s\n")

    overall_time = time_module.time() - overall_start
    print(f"\n{'='*60}")
    print(f"🎉 Presentation rebuild complete!")
    print(f"⏱️  Total time: {overall_time:.2f}s ({overall_time/60:.1f} minutes)")
    print(f"📊 Average per slide: {overall_time/total_slides:.2f}s")
    print(f"{'='*60}\n")
    
    return prs