
import streamlit as st
import os
import time
import base64
import io
import uuid
import csv
from datetime import datetime
from dotenv import load_dotenv

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, UnidentifiedImageError
from io import BytesIO
import google.generativeai as genai
from google.generativeai.types import GenerationConfig

# Project imports
from pptx_rag_quizzer.utils import parse_powerpoint, rebuild_presentation_with_accessible_features
from pptx_rag_quizzer.rag_core import RAGCore
from pptx_rag_quizzer.image import Image as ImageProcessor
from models.models import Presentation as PresentationModel

# Load environment variables
load_dotenv()

# Save consent response (email only when opted in)
def save_consent_email(email: str, choice: str):
    try:
        file_path = "consent_responses.csv"
        file_exists = os.path.exists(file_path)
        with open(file_path, "a", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if not file_exists:
                writer.writerow(["timestamp_utc", "email", "choice"])  # header
            writer.writerow([datetime.utcnow().isoformat(), email, choice])
    except Exception as e:
        print(f"Error saving consent response: {e}")


def parse_powerpoint_file(file_bytes, file_name):
    """Parse PowerPoint file into our model format"""
    file_object = io.BytesIO(file_bytes)
    return parse_powerpoint(file_object, file_name)

def process_powerpoint_with_rag_enhanced(pptx_model, enhanced_collection_id, file_bytes):
    """Process PowerPoint with RAG-enhanced accessibility features using the new rebuild method"""
    
    file_object = io.BytesIO(file_bytes)
    # Rebuild presentation with accessible features
    prs = rebuild_presentation_with_accessible_features(pptx_model, file_object)
    
    return prs


def main():
    st.set_page_config(
        page_title="PowerPoint Accessibility Enhancer",
        page_icon="♿",
        layout="centered"
    )
    
    st.title("♿ PowerPoint Accessibility Enhancer")
    st.markdown("Transform your PowerPoint into an accessible masterpiece with RAG-enhanced AI descriptions")
    st.warning("AI-Generated Content: The following text is created by an AI model. It may contain errors, omissions, or outdated information. Please exercise caution and confirm details before use.")
    
    # Initialize session state
    if 'processing_stage' not in st.session_state:
        st.session_state.processing_stage = 'consent'
    if 'presentation_model' not in st.session_state:
        st.session_state.presentation_model = None
    if 'rag_core' not in st.session_state:
        st.session_state.rag_core = None
    if 'image_processor' not in st.session_state:
        st.session_state.image_processor = None
    if 'collection_id' not in st.session_state:
        st.session_state.collection_id = None
    if 'current_batch' not in st.session_state:
        st.session_state.current_batch = 0
    if 'batch_size' not in st.session_state:
        st.session_state.batch_size = 5
    if 'consent_completed' not in st.session_state:
        st.session_state.consent_completed = False
    if 'consent_choice' not in st.session_state:
        st.session_state.consent_choice = None
    if 'consent_email' not in st.session_state:
        st.session_state.consent_email = ""
    if 'new_presentation_model' not in st.session_state:
        st.session_state.new_presentation_model = None
    
    # Consent Gate
    if st.session_state.processing_stage == 'consent':
        st.header("Consent to Participate")
        st.markdown(
            """
            Dear Colleagues,

            You are invited to participate in this research on “Enhancing Accessibility in Higher Education through Experiential Learning Opportunities” by participating a one-minute survey later (via email), or participating in one focus group discussion session, or one-on-one interview if you cannot make the focus group session. Confidentiality will be reminded at the beginning of the session.

            If you choose to volunteer to participate in the study, what you share in the focus group or interview will be included in the research analysis. The data will be aggregated and no participants will be identified.

            If you choose to opt out of the study, no further action is needed. There is no penalty if you choose not to participate.

            You can simply check one of the following boxes to indicate your participation.
            """
        )


        consent_options = [
            "Yes - I agree to participate in this research project. I am 18 years of age or older.",
            "Yes, I agree. I have provided my email address, and I confirm that I am 18 years of age or older.",
            "No - I do not agree to participate in this research project. I am 18 years of age or older.",
            "No - I am not eligible to participate as I am under the age of 18.",
        ]
        consent_choice = st.radio("Please select one option to continue:", consent_options, index=None)

        email = ""
        if consent_choice == consent_options[0]:
            email = st.text_input("Email (for 1-minute survey)", value=st.session_state.consent_email, placeholder="name@domain.edu")

        if st.button("Continue", type="primary", use_container_width=True):
            if consent_choice is None:
                st.error("Please select an option to continue.")
            elif consent_choice == consent_options[0]:
                if not email.strip():
                    st.error("Please enter your email to continue.")
                else:
                    # Save and proceed
                    st.session_state.consent_choice = 'yes'
                    st.session_state.consent_email = email.strip()
                    save_consent_email(st.session_state.consent_email, st.session_state.consent_choice)
                    st.session_state.consent_completed = True
                    st.session_state.processing_stage = 'upload'
                    st.rerun()
            elif consent_choice == consent_options[2] or consent_choice == consent_options[1]:
                st.session_state.consent_choice = 'no'
                st.session_state.consent_completed = True
                st.session_state.processing_stage = 'upload'
                st.rerun()
            else:
                st.session_state.consent_choice = 'under_18'
                st.session_state.consent_completed = False
                st.session_state.processing_stage = 'blocked'
                st.rerun()

    # Stage 0: Blocked (under 18)
    elif st.session_state.processing_stage == 'blocked':
        st.header("Access Restricted")
        st.error("You are not eligible to participate as you are under the age of 18. Please close this page.")
        st.stop()

    # Stage 1: File Upload and Initial Processing
    elif st.session_state.processing_stage == 'upload':
        if not st.session_state.consent_completed:
            st.session_state.processing_stage = 'consent'
            st.rerun()
        st.header("📁 Step 1: Upload PowerPoint")
        
        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file (Note: The uploaded file is not collected, and any temporary cache will be cleared automatically when the app is closed.",
            type=['pptx'],
            help="Upload a .pptx file to enhance accessibility"
        )
        
        if uploaded_file is not None:
            if st.button("🚀 Process PowerPoint", type="primary", use_container_width=True):
                with st.spinner("Processing PowerPoint and building RAG collection..."):
                    try:
                        # Parse the PowerPoint file
                        file_bytes = uploaded_file.read()
                        presentation_model = parse_powerpoint_file(file_bytes, uploaded_file.name)
                        
                        # Initialize RAG core and create collection from text content
                        rag_core = RAGCore()
                        collection_id = rag_core.create_collection(presentation_model)
                        
                        # Initialize image processor
                        image_processor = ImageProcessor(rag_core)
                        
                        # Store in session state
                        st.session_state.presentation_model = presentation_model
                        st.session_state.rag_core = rag_core
                        st.session_state.image_processor = image_processor
                        st.session_state.collection_id = collection_id
                        st.session_state.uploaded_file_name = uploaded_file.name
                        st.session_state.file_bytes = file_bytes
                        
                        st.session_state.processing_stage = 'describe_images'
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"❌ Error processing PowerPoint: {str(e)}")
                        st.exception(e)

            if st.button("⚡ Quick Generate & Download (Skip Review)", use_container_width=True):
                with st.spinner("Generating accessible PowerPoint (skipping review) This may take a few minutes depending on the size of the presentation..."):
                    try:
                        # Parse the PowerPoint file
                        file_bytes = uploaded_file.read()
                        presentation_model = parse_powerpoint_file(file_bytes, uploaded_file.name)

                        # Initialize RAG core and create initial collection from text content
                        rag_core = RAGCore()
                        collection_id = rag_core.create_collection(presentation_model)

                        # Initialize image processor
                        image_processor = ImageProcessor(rag_core)

                        # Auto-generate descriptions for all images without going through review UI
                        for slide in presentation_model.slides:
                            for item in slide.items:
                                if item.type.value == 'image':
                                    if not item.content or item.content.lower() in ['none', 'null', '']:
                                        try:
                                            image_description = image_processor.describe_image(
                                                item.image_bytes,
                                                item.extension,
                                                item.slide_number,
                                                collection_id
                                            )
                                            if image_description and image_description.startswith("Description: "):
                                                image_description = image_description[len("Description: "):]
                                            if image_description and image_description != "None":
                                                item.content = image_description
                                            else:
                                                item.content = "No description available"
                                        except Exception as e:
                                            item.content = f"Error describing image: {e}"

                        # Rebuild collection including image descriptions
                        response = rag_core.remove_collection(collection_id)
                        if not response.get("success", False):
                            st.error(f"❌ Error removing collection: {response.get('message', 'Unknown error')}")
                            st.stop()
                        enhanced_collection_id = rag_core.create_collection(presentation_model)

                        # Process with enhanced RAG and write accessibility features
                        prs = process_powerpoint_with_rag_enhanced(presentation_model, enhanced_collection_id, file_bytes)

                        # Jump straight to download stage
                        st.session_state.new_presentation_model = prs
                        st.session_state.uploaded_file_name = uploaded_file.name
                        st.session_state.processing_stage = 'download'
                        st.rerun()

                    except Exception as e:
                        st.error(f"❌ Error generating accessible PowerPoint: {str(e)}")
                        st.exception(e)
    
    # Stage 2: Image Description and Batch Processing
    elif st.session_state.processing_stage == 'describe_images':
        st.header("🖼️ Step 2: Describe Images")
        
        # Back button
        if st.button("← Back to Upload", use_container_width=True):
            st.session_state.processing_stage = 'upload'
            st.rerun()
        
        presentation_model = st.session_state.presentation_model
        image_processor = st.session_state.image_processor
        collection_id = st.session_state.collection_id
        
        # Get all images from the presentation
        all_images = []
        for slide in presentation_model.slides:
            for item in slide.items:
                if item.type.value == 'image':
                    all_images.append(item)
        
        total_images = len(all_images)
        
        if total_images == 0:
            st.info("📝 No images found in the presentation.")
            st.session_state.processing_stage = 'final_processing'
            st.rerun()
            return
        
        # Batch processing
        batch_size = st.session_state.batch_size
        current_batch = st.session_state.current_batch
        batch_start = current_batch * batch_size
        batch_end = min(batch_start + batch_size, total_images)
        current_batch_images = all_images[batch_start:batch_end]
        
        # Show progress
        st.write(f"**Processing images {batch_start + 1} to {batch_end} of {total_images}**")
        st.progress((batch_end) / total_images)
        
        # Process current batch if descriptions are missing
        batch_ready = all(
            img.content and img.content.lower() not in ['none', 'null', ''] 
            for img in current_batch_images
        )
        
        if not batch_ready:
            st.write("🤖 Generating AI descriptions for images...")
            
            with st.spinner("AI is analyzing images and generating descriptions..."):
                for i, img_item in enumerate(current_batch_images):
                    if not img_item.content or img_item.content.lower() in ['none', 'null', '']:
                        try:
                            st.write(f"Describing image {batch_start + i + 1} of {total_images}...")
                            image_description = image_processor.describe_image(
                                img_item.image_bytes,
                                img_item.extension,
                                img_item.slide_number,
                                collection_id
                            )
                            
                            # Clean up description
                            if image_description and image_description.startswith("Description: "):
                                image_description = image_description[len("Description: "):]
                            
                            if image_description and image_description != "None":
                                img_item.content = image_description
                                st.write(f"✓ Image {batch_start + i + 1} described successfully")
                            else:
                                img_item.content = "No description available"
                                st.write(f"⚠️ No description generated for image {batch_start + i + 1}")
                                
                        except Exception as e:
                            img_item.content = f"Error describing image: {e}"
                            st.write(f"✗ Error describing image {batch_start + i + 1}: {e}")
            
            st.success("All descriptions generated! Please review and approve.")
            st.rerun()
        
        else:
            # Display current batch for review
            st.write("📋 **Review and approve image descriptions:**")
            
            for i, img_item in enumerate(current_batch_images):
                st.write(f"**Image {batch_start + i + 1}** (Slide {img_item.slide_number})")
                
                # Display image
                #img = safe_open_image(img_item.image_bytes)
                img = Image.open(io.BytesIO(img_item.image_bytes))
                st.image(img, width=400, caption=f"Slide {img_item.slide_number} - Image {batch_start + i + 1}")
                
                # Editable description
                new_description = st.text_area(
                    f"Description for image {batch_start + i + 1}:",
                    value=img_item.content,
                    height=100,
                    key=f"desc_{batch_start + i}"
                )
                img_item.content = new_description
                st.write("---")
            
            # Navigation buttons
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if batch_start > 0:
                    if st.button("← Previous Batch"):
                        st.session_state.current_batch = current_batch - 1
                        st.rerun()
            
            with col2:
                if st.button("💾 Save Batch"):
                    st.success("Batch saved!")
            
            with col3:
                if batch_end < total_images:
                    if st.button("Next Batch →"):
                        st.session_state.current_batch = current_batch + 1
                        st.rerun()
                else:
                    # All images processed
                    if st.button("✅ Finish Image Processing", type="primary"):
                        st.session_state.processing_stage = 'final_processing'
                        st.rerun()
    
    # Stage 3: Final Processing with Enhanced RAG
    elif st.session_state.processing_stage == 'final_processing':
        st.header("🎯 Step 3: Final Processing")
        
        # Back button
        if st.button("← Back to Image Processing", use_container_width=True):
            st.session_state.processing_stage = 'describe_images'
            st.rerun()
        
        if st.button("🚀 Generate Enhanced Accessibility Features", type="primary", use_container_width=True):
            with st.spinner("Generating accessibility features... This may take a few minutes depending on the size of the presentation..."):
                try:
                    presentation_model = st.session_state.presentation_model
                    rag_core = st.session_state.rag_core
                    collection_id = st.session_state.collection_id
                    file_bytes = st.session_state.file_bytes
                    file_name = st.session_state.uploaded_file_name
                    
                    # Remove old collection and create new one with image descriptions
                    response = rag_core.remove_collection(collection_id)
                    if not response.get("success", False):
                        st.error(f"❌ Error removing collection: {response.get('message', 'Unknown error')}")
                        st.stop()
                    enhanced_collection_id = rag_core.create_collection(presentation_model)
                    
                    # Process with enhanced RAG
                    prs = process_powerpoint_with_rag_enhanced(presentation_model, enhanced_collection_id, file_bytes)
                    st.session_state.new_presentation_model = prs
                    
                    # Store output path for download
                    st.session_state.uploaded_file_name = file_name
                    st.session_state.processing_stage = 'download'
                    st.rerun()
                    
                except Exception as e:
                    st.error(f"❌ Error in final processing: {str(e)}")
                    st.exception(e)
    
    # Stage 4: Download
    elif st.session_state.processing_stage == 'download':
        prs = st.session_state.new_presentation_model
        file_name = st.session_state.uploaded_file_name
        output_path = f"accessible_{file_name}"
        
        # Save the presentation to a file
        prs.save(output_path)
        
        st.header("✅ Accessibility Enhancement Complete!")
        
        # Show results
        col1, col2 = st.columns(2)
        with col1:
            st.success("♿ Accessibility features added:")
            st.write("• **Enhanced Alt Text**: AI-generated descriptions with context")
            st.write("• **Comprehensive Slide Notes**: RAG-enhanced content summaries")
            st.write("• **Chart/Table Labels**: Accessibility labels for data elements")
            st.write("• **Context-Aware Descriptions**: Using presentation-wide knowledge")
        
        with col2:
            st.info("📊 Processing Summary:")
            st.write("• **RAG Collection**: Built with text and image content")
            st.write("• **Image Processing**: Enhanced with contextual descriptions")
            st.write("• **Notes Generation**: Context-retrieved and AI-enhanced")
        
        # Download button
        if os.path.exists(output_path):
            with open(output_path, "rb") as f:
                st.download_button(
                    label="📥 Download Accessible PowerPoint",
                    data=f.read(),
                    file_name=output_path,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                    use_container_width=True
                )
        
        # Reset button
        if st.button("🔄 Process Another Presentation", use_container_width=True):
            # Clean up
            if os.path.exists(output_path):
                os.remove(output_path)
            
            # Reset session state
            for key in ['processing_stage', 'presentation_model', 'rag_core', 'image_processor', 
                       'collection_id', 'current_batch', 'output_path']:
                if key in st.session_state:
                    del st.session_state[key]
            
            st.session_state.processing_stage = 'upload'
            st.rerun()

if __name__ == "__main__":
    main()
